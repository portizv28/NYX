from pptx import Presentation


def crear_presentacion(
    titulo="Presentación NYX",
    nombre="NYX_presentacion.pptx"
):

    prs = Presentation()


    diapositiva = prs.slides.add_slide(
        prs.slide_layouts[1]
    )


    diapositiva.shapes.title.text = titulo


    prs.save(nombre)


    return nombre